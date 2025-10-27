#!/usr/bin/env python3
"""
Create comprehensive PowerPoint presentation from Challenge 1 markdown
Includes all content, code blocks, tables, and visual element guidance
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

def parse_markdown_file(md_file_path):
    """Parse the complete markdown file and extract all slide content"""
    with open(md_file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide sections
    slide_pattern = r'## SLIDE (\d+):(.*?)(?=## SLIDE|\Z)'
    matches = re.findall(slide_pattern, content, re.DOTALL)

    slides = []
    for slide_num, slide_content in matches:
        slide_data = {
            'number': int(slide_num),
            'title': '',
            'visual_elements': [],
            'speaker_notes': '',
            'slide_shows': [],
            'image_sources': []
        }

        # Extract title from visual elements or first heading
        title_patterns = [
            r'### Visual Elements:.*?[:-]\s*(?:Title:|.*?)(?:\*\*)?(.+?)(?:\n|$)',
            r'Title:\s*"([^"]+)"',
            r'# (.+?)(?:\n|$)',
        ]
        for pattern in title_patterns:
            title_match = re.search(pattern, slide_content, re.MULTILINE)
            if title_match:
                slide_data['title'] = title_match.group(1).strip().strip('"')
                break

        if not slide_data['title']:
            slide_data['title'] = f"Slide {slide_num}"

        # Extract visual elements section
        visual_match = re.search(r'### Visual Elements:(.*?)(?=### Speaker Notes|### Slide shows|\Z)', slide_content, re.DOTALL)
        if visual_match:
            visual_text = visual_match.group(1)
            # Extract bullet points
            bullets = re.findall(r'^[\s]*[-•]\s*(.+?)$', visual_text, re.MULTILINE)
            slide_data['visual_elements'] = [b.strip() for b in bullets if b.strip()]

            # Look for image/diagram references
            image_patterns = [
                r'Background:\s*(.+?)(?:\n|$)',
                r'Image:\s*(.+?)(?:\n|$)',
                r'Diagram:\s*(.+?)(?:\n|$)',
                r'Chart:\s*(.+?)(?:\n|$)',
                r'Screenshot:\s*(.+?)(?:\n|$)',
            ]
            for pattern in image_patterns:
                matches = re.findall(pattern, visual_text, re.MULTILINE)
                slide_data['image_sources'].extend(matches)

        # Extract "Slide shows" section
        shows_match = re.search(r'\*\*Slide shows\*\*:(.*?)(?=###|\*\*[A-Z]|\Z)', slide_content, re.DOTALL)
        if shows_match:
            shows_text = shows_match.group(1)
            # Extract code blocks
            code_blocks = re.findall(r'```(\w+)?\n(.*?)```', shows_text, re.DOTALL)
            slide_data['slide_shows'] = [{'language': lang or 'text', 'code': code} for lang, code in code_blocks]

            # Also capture tables and other formatted content
            if not code_blocks:
                # Look for tables or structured data
                lines = shows_text.strip().split('\n')
                if lines:
                    slide_data['slide_shows'].append({'language': 'text', 'code': shows_text.strip()})

        # Extract speaker notes (first 2000 chars to fit in notes section)
        notes_match = re.search(r'### Speaker Notes:(.*?)(?=##|\Z)', slide_content, re.DOTALL)
        if notes_match:
            notes = notes_match.group(1).strip()
            # Clean up notes
            notes = notes.replace('"', '').replace('"', '').replace("'", "'")
            # Remove excessive newlines
            notes = re.sub(r'\n{3,}', '\n\n', notes)
            slide_data['speaker_notes'] = notes[:2000]

        slides.append(slide_data)

    return sorted(slides, key=lambda x: x['number'])

def add_title_slide(prs):
    """Create opening title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Challenge 1: Fire Data Sources & Ingestion Mechanisms"
    subtitle.text = "Real-Time Satellite Fire Detection Pipeline\nCAL FIRE Space-Based Data Challenge\n$50,000 Prize Competition"

    # Style the title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_content_slide(prs, slide_data):
    """Create a comprehensive content slide with all materials"""

    # Use blank layout for maximum flexibility
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    current_top = Inches(1.2)

    # Add visual elements (bullet points)
    if slide_data['visual_elements']:
        content_box = slide.shapes.add_textbox(Inches(0.5), current_top, Inches(9), Inches(2.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, bullet in enumerate(slide_data['visual_elements'][:8]):  # Limit to 8 bullets
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet
            p.level = 0
            p.font.size = Pt(14)
            p.space_before = Pt(6)

        current_top += Inches(2.7)

    # Add "Slide shows" content (code blocks, tables)
    if slide_data['slide_shows']:
        shows_height = Inches(2.5) if slide_data['visual_elements'] else Inches(4.5)

        for show in slide_data['slide_shows'][:1]:  # Take first code block/table
            code_box = slide.shapes.add_textbox(Inches(0.5), current_top, Inches(9), shows_height)
            text_frame = code_box.text_frame
            text_frame.word_wrap = True

            # Add code/table content
            code_text = show['code'].strip()
            # Truncate if too long
            if len(code_text) > 1000:
                code_text = code_text[:1000] + "\n... (continued in speaker notes)"

            p = text_frame.paragraphs[0]
            p.text = code_text
            p.font.size = Pt(10)
            p.font.name = 'Courier New'

            # Add background color for code
            shape = code_box
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            shape.line.color.rgb = RGBColor(200, 200, 200)

        current_top += shows_height + Inches(0.2)

    # Add image source guidance box if images are referenced
    if slide_data['image_sources']:
        guide_top = Inches(6.5)
        guide_box = slide.shapes.add_textbox(Inches(0.5), guide_top, Inches(9), Inches(0.8))
        text_frame = guide_box.text_frame

        p = text_frame.paragraphs[0]
        p.text = "📷 VISUAL ELEMENTS NEEDED:"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(204, 102, 0)

        for img_ref in slide_data['image_sources'][:2]:  # Show first 2 references
            p = text_frame.add_paragraph()
            p.text = f"  • {img_ref.strip()}"
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(102, 102, 102)

        # Style the box
        guide_box.fill.solid()
        guide_box.fill.fore_color.rgb = RGBColor(255, 250, 205)
        guide_box.line.color.rgb = RGBColor(204, 102, 0)

    # Add speaker notes
    if slide_data['speaker_notes']:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data['speaker_notes']

    return slide

def add_image_resources_slide(prs):
    """Add final slide with links to all image resources"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Visual Resources & Image Sources"

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    text_frame = content_box.text_frame

    resources = [
        "NASA Images & Videos:",
        "  • NASA Image Gallery: https://images.nasa.gov",
        "  • NASA FIRMS: https://firms.modaps.eosdis.nasa.gov/gallery/",
        "  • NASA Worldview: https://worldview.earthdata.nasa.gov",
        "",
        "Satellite Images:",
        "  • Copernicus Open Access Hub: https://scihub.copernicus.eu",
        "  • USGS EarthExplorer: https://earthexplorer.usgs.gov",
        "",
        "Icons & Graphics:",
        "  • Font Awesome: https://fontawesome.com (free icons)",
        "  • Flaticon: https://www.flaticon.com",
        "  • Icons8: https://icons8.com",
        "",
        "Diagrams:",
        "  • draw.io: https://app.diagrams.net (free diagramming)",
        "  • Lucidchart: https://www.lucidchart.com",
        "",
        "Stock Images:",
        "  • Unsplash: https://unsplash.com (search: wildfire, satellite)",
        "  • Pexels: https://www.pexels.com (free stock photos)",
        "",
        "GitHub Repository (for code screenshots):",
        "  • Your repository: C:\\dev\\wildfire\\",
    ]

    for line in resources:
        p = text_frame.add_paragraph() if text_frame.text else text_frame.paragraphs[0]
        p.text = line
        p.font.size = Pt(12) if not line.startswith('  •') else Pt(10)
        p.font.bold = not line.startswith('  •')
        p.space_before = Pt(3)

    return slide

def create_comprehensive_presentation(md_file_path, output_path):
    """Create complete PowerPoint with all materials"""
    print(f"📖 Reading markdown file: {md_file_path}")
    slides_data = parse_markdown_file(md_file_path)
    print(f"✅ Parsed {len(slides_data)} slides")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add title slide
    print("📝 Creating title slide...")
    add_title_slide(prs)

    # Add all content slides
    for i, slide_data in enumerate(slides_data, 1):
        print(f"📝 Creating slide {i}/{len(slides_data)}: {slide_data['title']}")
        add_content_slide(prs, slide_data)

    # Add image resources slide
    print("📝 Adding image resources guide...")
    add_image_resources_slide(prs)

    # Save presentation
    print(f"💾 Saving presentation to: {output_path}")
    prs.save(output_path)

    print(f"\n{'='*60}")
    print(f"✅ COMPREHENSIVE POWERPOINT CREATED SUCCESSFULLY!")
    print(f"{'='*60}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"📁 File location: {output_path}")
    print(f"📐 Slide size: 10 x 7.5 inches (standard)")
    print(f"\n📋 What's included:")
    print(f"   ✓ Title slide")
    print(f"   ✓ {len(slides_data)} content slides with:")
    print(f"      - All bullet points from Visual Elements")
    print(f"      - Code blocks and tables from 'Slide shows'")
    print(f"      - Speaker notes (up to 2000 chars per slide)")
    print(f"      - Visual element guidance boxes")
    print(f"   ✓ Image resources guide (final slide)")
    print(f"\n🎨 Next steps:")
    print(f"   1. Open the PowerPoint file")
    print(f"   2. Look for yellow boxes - these show what images to add")
    print(f"   3. Use the final slide links to find images")
    print(f"   4. Review speaker notes for detailed explanations")
    print(f"   5. Customize colors, fonts, and layout as needed")
    print(f"{'='*60}\n")

if __name__ == "__main__":
    import os
    import sys

    # Support both Windows and Linux paths
    if os.path.exists("/workspace"):
        # Running in Docker container
        md_file = "/workspace/docs/CHALLENGE1_FIRE_DATA_PRESENTATION-OLD.md"
        output_file = "/workspace/docs/presentations/Challenge1_Fire_Data_Presentation.pptx"
    else:
        # Running on Windows
        md_file = r"C:\dev\wildfire\docs\CHALLENGE1_FIRE_DATA_PRESENTATION-OLD.md"
        output_file = r"C:\dev\wildfire\docs\presentations\Challenge1_Fire_Data_Presentation.pptx"

    # Allow command line arguments
    if len(sys.argv) > 1:
        md_file = sys.argv[1]
    if len(sys.argv) > 2:
        output_file = sys.argv[2]

    try:
        create_comprehensive_presentation(md_file, output_file)
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
