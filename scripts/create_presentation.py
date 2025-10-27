#!/usr/bin/env python3
"""
Convert Challenge 1 markdown presentation to PowerPoint (.pptx) format
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def parse_markdown_slides(md_file_path):
    """Parse markdown file and extract slide content"""
    with open(md_file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide markers (## SLIDE)
    slide_pattern = r'## SLIDE (\d+):(.*?)(?=## SLIDE|\Z)'
    matches = re.findall(slide_pattern, content, re.DOTALL)

    slides = []
    for slide_num, slide_content in matches:
        # Extract title
        title_match = re.search(r'### Visual Elements:.*?-\s*(?:Title:|.*?)(?:\*\*)?([^\n]+)', slide_content)
        title = title_match.group(1).strip() if title_match else f"Slide {slide_num}"

        # Extract speaker notes
        notes_match = re.search(r'### Speaker Notes:(.*?)(?=###|\Z)', slide_content, re.DOTALL)
        notes = notes_match.group(1).strip() if notes_match else ""

        # Extract visual elements
        visual_match = re.search(r'### Visual Elements:(.*?)(?=### Speaker Notes)', slide_content, re.DOTALL)
        visual = visual_match.group(1).strip() if visual_match else ""

        slides.append({
            'number': int(slide_num),
            'title': title,
            'visual': visual,
            'notes': notes
        })

    return sorted(slides, key=lambda x: x['number'])

def create_title_slide(prs, title_text):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Challenge 1: Fire Data Sources & Ingestion Mechanisms"
    subtitle.text = "Real-Time Satellite Fire Detection Pipeline\nCAL FIRE Space-Based Data Challenge"

    return slide

def create_content_slide(prs, slide_data):
    """Create a content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout

    # Set title
    title = slide.shapes.title
    title.text = slide_data['title']

    # Add content box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    # Parse visual elements and add as bullet points
    visual_lines = [line.strip() for line in slide_data['visual'].split('\n') if line.strip() and line.strip().startswith('-')]

    for line in visual_lines[:10]:  # Limit to 10 bullets to avoid overcrowding
        clean_line = line.lstrip('- ').strip()
        p = text_frame.add_paragraph()
        p.text = clean_line
        p.level = 0
        p.font.size = Pt(18)

    # Add speaker notes
    if slide_data['notes']:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        # Clean up notes - take first 500 chars to avoid overflow
        clean_notes = slide_data['notes'][:1000].replace('"', '').replace('"', '')
        text_frame.text = clean_notes

    return slide

def create_presentation(md_file_path, output_path):
    """Create PowerPoint presentation from markdown"""
    print(f"Parsing markdown file: {md_file_path}")
    slides = parse_markdown_slides(md_file_path)
    print(f"Found {len(slides)} slides")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create title slide
    print("Creating title slide...")
    create_title_slide(prs, "Challenge 1: Fire Data Sources & Ingestion")

    # Create content slides
    for i, slide_data in enumerate(slides, 1):
        print(f"Creating slide {i}/{len(slides)}: {slide_data['title']}")
        create_content_slide(prs, slide_data)

    # Save presentation
    print(f"Saving presentation to: {output_path}")
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created successfully!")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   File: {output_path}")

if __name__ == "__main__":
    import os
    import sys

    # Support both Windows and Linux paths
    if os.path.exists("/workspace"):
        # Running in Docker container
        md_file = "/workspace/docs/CHALLENGE1_FIRE_DATA_PRESENTATION-OLD.md"
        output_file = "/workspace/docs/presentations/Challenge1_Fire_Data_Presentation.pptx"
    else:
        # Running on Windows
        md_file = r"C:\dev\wildfire\docs\CHALLENGE1_FIRE_DATA_PRESENTATION-OLD.md"
        output_file = r"C:\dev\wildfire\docs\presentations\Challenge1_Fire_Data_Presentation.pptx"

    # Allow command line arguments
    if len(sys.argv) > 1:
        md_file = sys.argv[1]
    if len(sys.argv) > 2:
        output_file = sys.argv[2]

    try:
        create_presentation(md_file, output_file)
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
